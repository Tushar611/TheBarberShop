from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Brand Colors
bg_dark = RGBColor(10, 10, 10)
text_main = RGBColor(245, 245, 245)
gold_primary = RGBColor(212, 175, 55)

# Helper function to apply styling to slides
def style_slide(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_dark

# Slide 1: Title
title_slide_layout = prs.slide_layouts[0]
slide1 = prs.slides.add_slide(title_slide_layout)
style_slide(slide1)
title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "The Barber Shop - Digital Transformation"
title.text_frame.paragraphs[0].font.color.rgb = gold_primary
title.text_frame.paragraphs[0].font.name = "Playfair Display"
title.text_frame.paragraphs[0].font.size = Pt(44)

subtitle.text = "Elevating a Local Business into a Luxury Brand\nBahadurgarh, Haryana"
subtitle.text_frame.paragraphs[0].font.color.rgb = text_main
subtitle.text_frame.paragraphs[0].font.name = "DM Sans"

# Slide 2: The Core Idea & Vision
bullet_slide_layout = prs.slide_layouts[1]
slide2 = prs.slides.add_slide(bullet_slide_layout)
style_slide(slide2)
title2 = slide2.shapes.title
body2 = slide2.placeholders[1]

title2.text = "The Core Idea & Vision"
title2.text_frame.paragraphs[0].font.color.rgb = gold_primary

tf2 = body2.text_frame
tf2.text = "Why build a website for a local Barber Shop?"
p = tf2.add_paragraph()
p.text = "To move away from being just another 'local shop' and establish a £20 Lakh premium, luxury identity."
p = tf2.add_paragraph()
p.text = "To automate the booking process so barbers can focus on cutting hair, not answering phone calls."
p = tf2.add_paragraph()
p.text = "To build instant trust with new clients searching on Google before they even step foot inside."
p = tf2.add_paragraph()
p.text = "To justify premium pricing through a high-end digital experience."

for para in tf2.paragraphs:
    para.font.color.rgb = text_main
    para.font.size = Pt(24)

# Slide 3: Key Services Showcased
slide3 = prs.slides.add_slide(bullet_slide_layout)
style_slide(slide3)
title3 = slide3.shapes.title
body3 = slide3.placeholders[1]

title3.text = "Premium Services Highlighted"
title3.text_frame.paragraphs[0].font.color.rgb = gold_primary

tf3 = body3.text_frame
tf3.text = "We built dedicated flows for the shop's core offerings:"
tf3.add_paragraph().text = "Modern Fade Cuts: Precision fade with sharp lines."
tf3.add_paragraph().text = "Beard Styling: Hot towel prep and razor sharp lines."
tf3.add_paragraph().text = "Classic Haircuts & Combos: Traditional shears cut and wash."
tf3.add_paragraph().text = "Kids Haircuts & General Grooming"
tf3.add_paragraph().text = "Every service has a direct 'Book Now' button that auto-selects the service in the booking form."

for para in tf3.paragraphs:
    para.font.color.rgb = text_main
    para.font.size = Pt(24)

# Slide 4: Major Features & UI/UX Upgrades (Part 1)
slide4 = prs.slides.add_slide(bullet_slide_layout)
style_slide(slide4)
title4 = slide4.shapes.title
body4 = slide4.placeholders[1]

title4.text = "Major Website Features & Upgrades"
title4.text_frame.paragraphs[0].font.color.rgb = gold_primary

tf4 = body4.text_frame
tf4.text = "Dynamic WhatsApp Booking Flow: Forms send structured info directly to the shop's WhatsApp (+91 8709298368). Includes a 'Next available slot' urgency driver."
tf4.add_paragraph().text = "Cinematic Parallax Hero: A full-screen, highly realistic luxury background that wows the user immediately."
tf4.add_paragraph().text = "Luxury Aesthetics: Dark noise textures, gold-fill interactive buttons, Playfair Display headers, and custom gold dividers."
tf4.add_paragraph().text = "Trust & Social Proof: '500+ Happy Clients', '5.0 Google Rating', and direct links to actual Google Reviews."

for para in tf4.paragraphs:
    para.font.color.rgb = text_main
    para.font.size = Pt(20)

# Slide 5: Major Features & UI/UX Upgrades (Part 2)
slide5 = prs.slides.add_slide(bullet_slide_layout)
style_slide(slide5)
title5 = slide5.shapes.title
body5 = slide5.placeholders[1]

title5.text = "Technical & Mobile Features"
title5.text_frame.paragraphs[0].font.color.rgb = gold_primary

tf5 = body5.text_frame
tf5.text = "100% Mobile Optimized: Includes a modern dark slide-in drawer for mobile navigation, replacing basic menus."
tf5.add_paragraph().text = "Sticky Glassmorphism Header: The 'Book Now' button follows the user as they scroll, capturing impulse bookings."
tf5.add_paragraph().text = "Interactive Google Map: Embedded directly on the contact page for seamless navigation to Shop No 53, Omaxe Galleria."
tf5.add_paragraph().text = "Page Loader & Animations: Custom gold spinner on initial load and smooth scroll-triggered fade-ins."

for para in tf5.paragraphs:
    para.font.color.rgb = text_main
    para.font.size = Pt(22)

# Slide 6: The Bottom Line Benefits
slide6 = prs.slides.add_slide(bullet_slide_layout)
style_slide(slide6)
title6 = slide6.shapes.title
body6 = slide6.placeholders[1]

title6.text = "The Bottom Line (Why this matters)"
title6.text_frame.paragraphs[0].font.color.rgb = gold_primary

tf6 = body6.text_frame
tf6.text = "1. 24/7 Virtual Receptionist (Bookings happen while you sleep)"
tf6.add_paragraph().text = "2. Immediate Authority (You look like the undisputed best in town)"
tf6.add_paragraph().text = "3. Zero Location Confusion (Integrated maps guide them directly to the chair)"
tf6.add_paragraph().text = "4. Higher Google Ranking (SEO optimization pushes you above competitors)"
tf6.add_paragraph().text = "5. Increased Revenue (Premium brand perception allows for premium pricing)"

for para in tf6.paragraphs:
    para.font.color.rgb = text_main
    para.font.size = Pt(24)

# Save presentation
prs.save('d:/websites/Barber_Shop_Pitch.pptx')
print("Successfully created presentation at d:/websites/Barber_Shop_Pitch.pptx")
